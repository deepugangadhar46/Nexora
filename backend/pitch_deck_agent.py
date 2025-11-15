"""
PITCH DECK AGENT for Nexora AI IDE
===================================

A professional AI agent that creates investor-ready pitch decks automatically using:
- Groq API (Llama) for content generation
- ElevenLabs API for voiceover narration
- QuickChart API for chart generation
- Auto Slide Generator (Problem → Solution → Market → Team → Ask)
- AI Design Theme Selector
- Script Writer for Demo Day
- Slide Export as PPTX/Video
- Investor Q&A Simulator

Author: NEXORA Team
Version: 1.0.0
License: MIT
"""

import os
import json
import uuid
import asyncio
import logging
import re
import base64
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass, asdict, field
from io import BytesIO
from urllib.parse import quote

# Third-party imports
import aiohttp
import requests
from dotenv import load_dotenv

# PPTX Generation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PPTX generation will be disabled.")

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


# ============================================================================
# DATA CLASSES
# ============================================================================

@dataclass
class SlideContent:
    """Content for a single slide"""
    slide_number: int
    title: str
    content: List[str]
    notes: str = ""
    chart_data: Optional[Dict[str, Any]] = None
    chart_type: Optional[str] = None  # bar, line, pie, doughnut


@dataclass
class PitchDeckSlides:
    """Complete pitch deck slides"""
    title_slide: SlideContent
    problem_slide: SlideContent
    solution_slide: SlideContent
    market_slide: SlideContent
    product_slide: SlideContent
    business_model_slide: SlideContent
    traction_slide: SlideContent
    competition_slide: SlideContent
    team_slide: SlideContent
    financials_slide: SlideContent
    ask_slide: SlideContent
    closing_slide: SlideContent


@dataclass
class DesignTheme:
    """Design theme for pitch deck"""
    name: str
    primary_color: str
    secondary_color: str
    accent_color: str
    background_color: str
    text_color: str
    font_family: str
    style_description: str


@dataclass
class VoiceoverNarration:
    """Voiceover narration for slides"""
    slide_number: int
    text: str
    audio_url: Optional[str] = None
    duration_seconds: float = 0.0


@dataclass
class DemoScript:
    """Demo day pitch script"""
    full_script: str
    slide_scripts: List[Dict[str, Any]]
    total_duration_minutes: float
    pacing_cues: List[str]
    emphasis_points: List[str]


@dataclass
class InvestorQuestion:
    """Investor Q&A simulation"""
    question: str
    category: str  # financial, market, team, product, competition
    difficulty: str  # easy, medium, hard
    suggested_answer: str
    key_points: List[str]


@dataclass
class PitchDeckResponse:
    """Complete pitch deck response"""
    deck_id: str
    business_name: str
    tagline: str
    slides: PitchDeckSlides
    design_theme: DesignTheme
    voiceovers: List[VoiceoverNarration]
    demo_script: DemoScript
    investor_qa: List[InvestorQuestion]
    pptx_url: Optional[str] = None
    video_url: Optional[str] = None
    created_at: str = ""


# ============================================================================
# API CLIENTS
# ============================================================================

class GroqClient:
    """Groq API client for Llama models with rate limiting"""
    
    def __init__(self, api_key: Optional[str] = None):
        self.api_key = api_key or os.getenv("GROQ_API_KEY")
        self.base_url = "https://api.groq.com/openai/v1/chat/completions"
        # Use a concrete Groq-supported Llama model by default
        # You can still override via GROQ_MODEL if needed.
        self.model = os.getenv("GROQ_MODEL", "llama-3.1-8b-instant")
        self.last_request_time = 0
        self.min_request_interval = 2.0  # Minimum 2 seconds between requests
        
        if not self.api_key:
            raise ValueError("GROQ_API_KEY not found in environment variables")
        
        logger.info(f"GroqClient initialized with model: {self.model} (rate limited)")
    
    async def generate(
        self,
        prompt: str,
        system_prompt: Optional[str] = None,
        temperature: float = 0.7,
        max_tokens: int = 4000,
        json_mode: bool = False,
        model: Optional[str] = None,
        max_retries: int = 5
    ) -> str:
        """Generate response using Groq with exponential backoff retry"""
        
        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": prompt})
        
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": model or self.model,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_tokens,
            "top_p": 0.95,
            "frequency_penalty": 0.1,
            "presence_penalty": 0.1
        }
        
        if json_mode:
            payload["response_format"] = {"type": "json_object"}
        
        retry_count = 0
        base_delay = 2.0
        
        while retry_count <= max_retries:
            try:
                # Rate limiting: ensure minimum interval between requests
                import time
                current_time = time.time()
                time_since_last = current_time - self.last_request_time
                if time_since_last < self.min_request_interval:
                    wait_time = self.min_request_interval - time_since_last
                    logger.info(f"Rate limiting: waiting {wait_time:.1f}s before request")
                    await asyncio.sleep(wait_time)
                
                async with aiohttp.ClientSession() as session:
                    async with session.post(
                        self.base_url,
                        json=payload,
                        headers=headers,
                        timeout=aiohttp.ClientTimeout(total=60)
                    ) as response:
                        self.last_request_time = time.time()
                        
                        if response.status == 429:  # Rate limit error
                            retry_count += 1
                            if retry_count > max_retries:
                                error_text = await response.text()
                                raise Exception(f"Rate limit exceeded after {max_retries} retries. Please wait 60-90 seconds and try again. Error: {error_text}")
                            
                            # Exponential backoff: 2s, 4s, 8s, 16s, 32s
                            wait_time = base_delay * (2 ** retry_count)
                            logger.warning(f"Rate limit hit. Retry {retry_count}/{max_retries} after {wait_time}s...")
                            await asyncio.sleep(wait_time)
                            continue
                        
                        if response.status != 200:
                            error_text = await response.text()
                            logger.error(f"Groq API error: {error_text}")
                            raise Exception(f"Groq API error ({response.status}): {error_text}")
                        
                        data = await response.json()
                        return data["choices"][0]["message"]["content"]
                        
            except asyncio.TimeoutError:
                retry_count += 1
                if retry_count > max_retries:
                    raise Exception(f"Request timeout after {max_retries} retries")
                wait_time = base_delay * retry_count
                logger.warning(f"Timeout. Retry {retry_count}/{max_retries} after {wait_time}s...")
                await asyncio.sleep(wait_time)
                
            except Exception as e:
                if "rate limit" in str(e).lower() or "429" in str(e):
                    retry_count += 1
                    if retry_count > max_retries:
                        raise Exception(f"Rate limit exceeded. Please wait 60-90 seconds and try again. The free Groq API has a limit of 12,000 tokens per minute.")
                    wait_time = base_delay * (2 ** retry_count)
                    logger.warning(f"Rate limit error. Retry {retry_count}/{max_retries} after {wait_time}s...")
                    await asyncio.sleep(wait_time)
                else:
                    logger.error(f"Error calling Groq API: {str(e)}")
                    raise
        
        raise Exception("Maximum retries exceeded")


class MoonshotClient:
    """Moonshot (Kimi) API client, OpenAI-compatible chat endpoint"""

    def __init__(self, api_key: Optional[str] = None):
        self.api_key = api_key or os.getenv("MOONSHOT_API_KEY")
        self.base_url = "https://api.moonshot.ai/v1/chat/completions"
        # Default Kimi model; adjust to the latest recommended model if needed
        self.model = os.getenv("MOONSHOT_MODEL", "kimi-k2-0711-preview")

        if not self.api_key:
            raise ValueError("MOONSHOT_API_KEY not found in environment variables")

    async def generate(
        self,
        prompt: str,
        system_prompt: Optional[str] = None,
        temperature: float = 0.7,
        max_tokens: int = 4000,
        json_mode: bool = False,
        model: Optional[str] = None
    ) -> str:
        """Generate response using Moonshot chat completions"""

        messages: List[Dict[str, str]] = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        messages.append({"role": "user", "content": prompt})

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
        }

        payload: Dict[str, Any] = {
            "model": model or self.model,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_tokens,
        }

        if json_mode:
            # Moonshot is OpenAI-compatible; use response_format if supported
            payload["response_format"] = {"type": "json_object"}

        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    self.base_url,
                    json=payload,
                    headers=headers,
                    timeout=aiohttp.ClientTimeout(total=60),
                ) as response:
                    if response.status != 200:
                        error_text = await response.text()
                        logger.error(f"Moonshot API error: {error_text}")
                        raise Exception(f"Moonshot API error: {response.status}")

                    data = await response.json()
                    return data["choices"][0]["message"]["content"]

        except Exception as e:
            logger.error(f"Error calling Moonshot API: {str(e)}")
            raise


class ElevenLabsClient:
    """ElevenLabs API client for voiceover narration"""
    
    def __init__(self, api_key: Optional[str] = None):
        self.api_key = api_key or os.getenv("ELEVENLABS_API_KEY")
        self.base_url = "https://api.elevenlabs.io/v1"
        self.default_voice_id = "21m00Tcm4TlvDq8ikWAM"  # Rachel voice
        
        if not self.api_key:
            logger.warning("ELEVENLABS_API_KEY not found. Voiceover will be disabled.")
    
    async def generate_speech(
        self,
        text: str,
        voice_id: Optional[str] = None,
        model_id: str = "eleven_monolingual_v1"
    ) -> Optional[bytes]:
        """Generate speech from text"""
        
        if not self.api_key:
            logger.warning("ElevenLabs API key not set, skipping speech generation")
            return None
        
        voice_id = voice_id or self.default_voice_id
        url = f"{self.base_url}/text-to-speech/{voice_id}"
        
        headers = {
            "Accept": "audio/mpeg",
            "Content-Type": "application/json",
            "xi-api-key": self.api_key
        }
        
        payload = {
            "text": text,
            "model_id": model_id,
            "voice_settings": {
                "stability": 0.5,
                "similarity_boost": 0.5
            }
        }
        
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    url,
                    json=payload,
                    headers=headers,
                    timeout=aiohttp.ClientTimeout(total=60)
                ) as response:
                    if response.status == 200:
                        return await response.read()
                    else:
                        error_text = await response.text()
                        logger.error(f"ElevenLabs API error: {error_text}")
                        return None
        
        except Exception as e:
            logger.error(f"Error calling ElevenLabs API: {str(e)}")
            return None


class QuickChartClient:
    """QuickChart API client for chart generation"""
    
    def __init__(self):
        self.base_url = "https://quickchart.io/chart"
    
    def generate_chart_url(
        self,
        chart_type: str,
        data: Dict[str, Any],
        width: int = 800,
        height: int = 400
    ) -> str:
        """Generate chart URL using QuickChart API"""
        
        chart_config = {
            "type": chart_type,
            "data": data,
            "options": {
                "responsive": True,
                "maintainAspectRatio": True,
                "plugins": {
                    "legend": {
                        "display": True,
                        "position": "top"
                    }
                }
            }
        }
        
        config_json = json.dumps(chart_config)
        encoded_config = quote(config_json)
        
        return f"{self.base_url}?width={width}&height={height}&chart={encoded_config}"
    
    async def generate_chart_image(
        self,
        chart_type: str,
        data: Dict[str, Any],
        width: int = 800,
        height: int = 400
    ) -> Optional[bytes]:
        """Generate chart image bytes"""
        
        url = self.generate_chart_url(chart_type, data, width, height)
        
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=aiohttp.ClientTimeout(total=30)) as response:
                    if response.status == 200:
                        return await response.read()
                    else:
                        logger.error(f"QuickChart API error: {response.status}")
                        return None
        
        except Exception as e:
            logger.error(f"Error calling QuickChart API: {str(e)}")
            return None


# ============================================================================
# PITCH DECK AGENT
# ============================================================================

class PitchDeckAgent:
    """
    PITCH DECK AGENT
    
    Creates professional, investor-ready pitch decks automatically with:
    - Auto Slide Generator (12 slides)
    - Voiceover Narration (ElevenLabs)
    - AI Design Theme Selector
    - Chart Auto-Builder (QuickChart)
    - Script Writer for Demo Day
    - Slide Export as PPTX/Video
    - Investor Q&A Simulator
    """
    
    def __init__(
        self,
        groq_api_key: Optional[str] = None,
        elevenlabs_api_key: Optional[str] = None
    ):
        """Initialize the Pitch Deck Agent"""
        
        # Always use Groq client for Llama models
        logger.info("Initializing PitchDeckAgent with Groq client")
        self.groq = GroqClient(groq_api_key)
        self.elevenlabs = ElevenLabsClient(elevenlabs_api_key)
        self.quickchart = QuickChartClient()
        
        logger.info("Pitch Deck Agent initialized successfully")
    
    # ========================================================================
    # MODULE 1: AUTO SLIDE GENERATOR
    # ========================================================================
    
    async def generate_slides(
        self,
        business_idea: str,
        business_name: str = "",
        target_market: str = "",
        funding_ask: float = 0
    ) -> PitchDeckSlides:
        """Generate structured slide content for the pitch deck.

        This uses the LLM to create 12 slides with clear titles and
        30 concise, investor-focused bullet points per slide.
        """
        
        system_prompt = """You are a world-class startup pitch deck expert.
Generate a structured, investor-ready 12-slide pitch deck.

The deck MUST follow this exact structure and order:
1. Title Slide - Company name, one-line tagline, what you do
2. Problem - 23 concrete pain points and who experiences them
3. Solution - Product/service overview and how it solves the problem
4. Market Opportunity - TAM / SAM / SOM and market context
5. Product - Key features and how users interact with it
6. Business Model - How you make money (pricing, revenue streams)
7. Traction & Metrics - Growth, users, revenue, key milestones
8. Competition - Competitors and your differentiation/moat
9. Team - Founders, roles, and relevant experience
10. Financials - High-level 35 year projections and unit economics
11. The Ask & Use of Funds - How much you are raising and where it goes
12. Closing / Vision - Long-term vision and call to action

Write content that is specific to the business described by the user.
Avoid generic phrases like "innovative solution" or "cutting-edge platform".
Use 36 concise bullet points per slide that an investor can read quickly.

Return ONLY valid JSON with this structure (no extra text):
{
  "slides": [
    {
      "slide_number": 1,
      "title": "Slide Title",
      "content": ["bullet point 1", "bullet point 2", "bullet point 3"],
      "notes": "Speaker notes for this slide"
    }
  ]
}"""

        funding_ask_str = f"${funding_ask:,.0f}" if funding_ask and funding_ask > 0 else "Not specified"

        user_prompt = f"""Create a complete 12-slide pitch deck for this business.

BUSINESS IDEA: {business_idea}
BUSINESS NAME: {business_name or "Not specified"}
TARGET MARKET: {target_market or "Not specified"}
FUNDING ASK: {funding_ask_str}

Requirements:
- Follow exactly the 12-slide structure defined in the system prompt.
- Each slide must have a clear, specific title.
- Each slide's "content" must be a list of 36 short, punchy bullet points.
- Make the content concrete, investor-focused, and easy to present.
- Where possible, include simple numbers or examples instead of vague claims.

Return ONLY valid JSON matching the specified schema."""

        try:
            response = await self.groq.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.6,
                json_mode=True,
                max_tokens=4000
            )

            data = json.loads(response)
            slides_data = data.get("slides", [])

            # Parse slides into structured format
            slides_dict: Dict[int, SlideContent] = {}
            for slide in slides_data:
                slides_dict[slide["slide_number"]] = SlideContent(
                    slide_number=slide["slide_number"],
                    title=slide["title"],
                    content=slide["content"],
                    notes=slide.get("notes", "")
                )

            # Ensure title slide is always customized with business name & idea
            title_slide = slides_dict.get(1)
            if not title_slide:
                # Create a title slide if the model failed to return one
                tagline = business_idea.split(".\n")[0].split(".")[0].strip()
                if not tagline:
                    tagline = business_idea[:120].strip()
                slides_dict[1] = SlideContent(
                    slide_number=1,
                    title=business_name or "Your Company",
                    content=[tagline] if tagline else ["AI-powered startup"],
                    notes="Auto-generated title slide based on business idea."
                )
            else:
                # Inject business name into the title if missing
                if business_name and business_name not in title_slide.title:
                    title_slide.title = business_name
                # Ensure first bullet acts as a tagline based on the idea
                if not title_slide.content:
                    tagline = business_idea.split(".\n")[0].split(".")[0].strip()
                    if not tagline:
                        tagline = business_idea[:120].strip()
                    title_slide.content = [tagline] if tagline else ["AI-powered startup"]
                slides_dict[1] = title_slide

            return PitchDeckSlides(
                title_slide=slides_dict.get(1, self._default_slide(1, "Title")),
                problem_slide=slides_dict.get(2, self._default_slide(2, "Problem")),
                solution_slide=slides_dict.get(3, self._default_slide(3, "Solution")),
                market_slide=slides_dict.get(4, self._default_slide(4, "Market")),
                product_slide=slides_dict.get(5, self._default_slide(5, "Product")),
                business_model_slide=slides_dict.get(6, self._default_slide(6, "Business Model")),
                traction_slide=slides_dict.get(7, self._default_slide(7, "Traction")),
                competition_slide=slides_dict.get(8, self._default_slide(8, "Competition")),
                team_slide=slides_dict.get(9, self._default_slide(9, "Team")),
                financials_slide=slides_dict.get(10, self._default_slide(10, "Financials")),
                ask_slide=slides_dict.get(11, self._default_slide(11, "The Ask")),
                closing_slide=slides_dict.get(12, self._default_slide(12, "Closing"))
            )
        
        except Exception as e:
            logger.error(f"Error generating slides: {str(e)}")
            return self._get_default_slides()
    
    # ========================================================================
    # MODULE 2: VOICEOVER NARRATION
    # ========================================================================
    
    async def generate_voiceovers(
        self,
        slides: PitchDeckSlides
    ) -> List[VoiceoverNarration]:
        """
        Add AI voiceover to slides using ElevenLabs
        """
        
        voiceovers = []
        all_slides = [
            slides.title_slide, slides.problem_slide, slides.solution_slide,
            slides.market_slide, slides.product_slide, slides.business_model_slide,
            slides.traction_slide, slides.competition_slide, slides.team_slide,
            slides.financials_slide, slides.ask_slide, slides.closing_slide
        ]
        
        for slide in all_slides:
            # Generate narration script
            narration_text = await self._generate_narration_script(slide)
            
            # Generate audio
            audio_bytes = await self.elevenlabs.generate_speech(narration_text)
            
            # Calculate duration (rough estimate: 150 words per minute)
            word_count = len(narration_text.split())
            duration = calculate_slide_duration(word_count)
            
            voiceover = VoiceoverNarration(
                slide_number=slide.slide_number,
                text=narration_text,
                audio_url=None,  # Would save to storage and return URL
                duration_seconds=duration
            )
            
            voiceovers.append(voiceover)
        
        return voiceovers
    
    async def _generate_narration_script(self, slide: SlideContent) -> str:
        """Generate natural narration script for a slide"""
        
        system_prompt = """You are a professional pitch coach. 
Convert slide bullet points into a natural, engaging narration script.
Keep it conversational, confident, and concise (30-60 seconds per slide)."""
        
        user_prompt = f"""Create a narration script for this slide:

TITLE: {slide.title}
CONTENT:
{chr(10).join(f"- {item}" for item in slide.content)}

NOTES: {slide.notes}

Write a natural, engaging script that sounds great when spoken aloud.
Keep it to 30-60 seconds. Be confident and compelling."""

        try:
            response = await self.groq.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.7,
                max_tokens=500
            )
            return response.strip()
        
        except Exception as e:
            logger.error(f"Error generating narration: {str(e)}")
            return f"{slide.title}. " + " ".join(slide.content)
    
    # ========================================================================
    # MODULE 3: AI DESIGN THEME SELECTOR
    # ========================================================================
    
    async def select_design_theme(
        self,
        business_idea: str,
        brand_tone: str = "professional"
    ) -> DesignTheme:
        """
        Create slide design based on brand tone (e.g., "Minimalist", "Techy")
        """
        
        system_prompt = """You are a UI/UX designer specializing in pitch deck design.
Select appropriate colors, fonts, and style for the given business.

Return ONLY valid JSON:
{
  "name": "Theme Name",
  "primary_color": "#HEX",
  "secondary_color": "#HEX",
  "accent_color": "#HEX",
  "background_color": "#HEX",
  "text_color": "#HEX",
  "font_family": "Font Name",
  "style_description": "Description of the design style"
}"""

        user_prompt = f"""Select a design theme for this pitch deck:

BUSINESS IDEA: {business_idea}
BRAND TONE: {brand_tone}

Choose colors and fonts that match the business and brand tone.
Popular tones: Minimalist, Techy, Bold, Corporate, Creative, Modern.
Return as JSON."""

        try:
            response = await self.groq.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.5,
                json_mode=True,
                max_tokens=500
            )
            
            data = json.loads(response)
            
            return DesignTheme(
                name=data.get("name", "Professional"),
                primary_color=data.get("primary_color", "#2563eb"),
                secondary_color=data.get("secondary_color", "#1e40af"),
                accent_color=data.get("accent_color", "#f59e0b"),
                background_color=data.get("background_color", "#ffffff"),
                text_color=data.get("text_color", "#1f2937"),
                font_family=data.get("font_family", "Inter, sans-serif"),
                style_description=data.get("style_description", "Clean and professional")
            )
        
        except Exception as e:
            logger.error(f"Error selecting design theme: {str(e)}")
            return self._get_default_theme()
    
    # ========================================================================
    # MODULE 4: CHART AUTO-BUILDER
    # ========================================================================
    
    async def add_charts_to_slides(
        self,
        slides: PitchDeckSlides,
        business_idea: str
    ) -> PitchDeckSlides:
        """
        Add charts for market size, growth, and revenue forecasts using QuickChart
        """
        
        # Generate market size chart for Market slide
        market_chart_data = await self._generate_market_chart_data(business_idea)
        slides.market_slide.chart_data = market_chart_data
        slides.market_slide.chart_type = "bar"
        
        # Generate revenue forecast chart for Financials slide
        revenue_chart_data = await self._generate_revenue_chart_data(business_idea)
        slides.financials_slide.chart_data = revenue_chart_data
        slides.financials_slide.chart_type = "line"
        
        # Generate traction chart for Traction slide
        traction_chart_data = await self._generate_traction_chart_data(business_idea)
        slides.traction_slide.chart_data = traction_chart_data
        slides.traction_slide.chart_type = "line"
        
        return slides
    
    async def _generate_market_chart_data(self, business_idea: str) -> Dict[str, Any]:
        """Generate market size chart data"""
        
        system_prompt = """You are a market research analyst.
Estimate TAM (Total Addressable Market), SAM (Serviceable Addressable Market), 
and SOM (Serviceable Obtainable Market) for the given business.

Return ONLY valid JSON:
{
  "labels": ["TAM", "SAM", "SOM"],
  "datasets": [{
    "label": "Market Size ($B)",
    "data": [100, 20, 2],
    "backgroundColor": ["#3b82f6", "#60a5fa", "#93c5fd"]
  }]
}"""

        user_prompt = f"""Estimate market sizes for: {business_idea}

Provide realistic TAM, SAM, SOM estimates in billions of dollars.
Return as JSON for a bar chart."""

        try:
            response = await self.groq.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.4,
                json_mode=True,
                max_tokens=500
            )
            return json.loads(response)
        
        except Exception as e:
            logger.error(f"Error generating market chart: {str(e)}")
            return {
                "labels": ["TAM", "SAM", "SOM"],
                "datasets": [{
                    "label": "Market Size ($B)",
                    "data": [50, 10, 1],
                    "backgroundColor": ["#3b82f6", "#60a5fa", "#93c5fd"]
                }]
            }
    
    async def _generate_revenue_chart_data(self, business_idea: str) -> Dict[str, Any]:
        """Generate revenue forecast chart data"""
        
        return {
            "labels": ["Year 1", "Year 2", "Year 3", "Year 4", "Year 5"],
            "datasets": [{
                "label": "Revenue ($M)",
                "data": [0.5, 2, 5, 12, 25],
                "borderColor": "#10b981",
                "backgroundColor": "rgba(16, 185, 129, 0.1)",
                "fill": True
            }]
        }
    
    async def _generate_traction_chart_data(self, business_idea: str) -> Dict[str, Any]:
        """Generate traction/growth chart data"""
        
        return {
            "labels": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            "datasets": [{
                "label": "Users",
                "data": [100, 250, 500, 1200, 2500, 5000],
                "borderColor": "#8b5cf6",
                "backgroundColor": "rgba(139, 92, 246, 0.1)",
                "fill": True
            }]
        }
    
    # ========================================================================
    # MODULE 5: SCRIPT WRITER FOR DEMO DAY
    # ========================================================================
    
    async def generate_demo_script(
        self,
        slides: PitchDeckSlides,
        target_duration_minutes: float = 5.0
    ) -> DemoScript:
        """
        Generate a 2-minute pitch script with pacing cues
        """
        
        system_prompt = """You are a pitch coach who has trained hundreds of founders for Demo Day.
Create a compelling, well-paced pitch script that flows naturally.

Include:
- Natural transitions between slides
- Pacing cues (PAUSE, EMPHASIZE, SLOW DOWN, etc.)
- Time allocation per slide
- Key emphasis points

Return ONLY valid JSON:
{
  "full_script": "Complete script text...",
  "slide_scripts": [
    {
      "slide_number": 1,
      "script": "Script for this slide...",
      "duration_seconds": 20,
      "pacing_cues": ["PAUSE after company name", "EMPHASIZE tagline"]
    }
  ],
  "total_duration_minutes": 5.0,
  "pacing_cues": ["Overall pacing tip 1", "Overall pacing tip 2"],
  "emphasis_points": ["Key point to emphasize 1", "Key point to emphasize 2"]
}"""

        # Collect all slide content
        all_slides = [
            slides.title_slide, slides.problem_slide, slides.solution_slide,
            slides.market_slide, slides.product_slide, slides.business_model_slide,
            slides.traction_slide, slides.competition_slide, slides.team_slide,
            slides.financials_slide, slides.ask_slide, slides.closing_slide
        ]
        
        slides_summary = "\n\n".join([
            f"SLIDE {slide.slide_number}: {slide.title}\n" + 
            "\n".join(f"- {item}" for item in slide.content)
            for slide in all_slides
        ])
        
        user_prompt = f"""Create a {target_duration_minutes}-minute Demo Day pitch script for these slides:

{slides_summary}

Requirements:
- Total duration: {target_duration_minutes} minutes
- Natural, conversational tone
- Smooth transitions
- Include pacing cues
- Emphasize key points
- Build excitement and momentum

Return as JSON."""

        try:
            response = await self.groq.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.7,
                json_mode=True,
                max_tokens=3000
            )
            
            data = json.loads(response)
            
            return DemoScript(
                full_script=data.get("full_script", ""),
                slide_scripts=data.get("slide_scripts", []),
                total_duration_minutes=float(data.get("total_duration_minutes", target_duration_minutes)),
                pacing_cues=data.get("pacing_cues", []),
                emphasis_points=data.get("emphasis_points", [])
            )
        
        except Exception as e:
            logger.error(f"Error generating demo script: {str(e)}")
            return self._get_default_demo_script()
    
    # ========================================================================
    # MODULE 6: INVESTOR Q&A SIMULATOR
    # ========================================================================
    
    async def generate_investor_qa(
        self,
        business_idea: str,
        slides: PitchDeckSlides,
        num_questions: int = 10
    ) -> List[InvestorQuestion]:
        """
        Let users practice answering questions via AI chat using AI reasoning
        """
        
        system_prompt = """You are a veteran venture capitalist conducting due diligence.
Generate tough but fair questions that investors would ask after a pitch.

Categories: financial, market, team, product, competition
Difficulty: easy, medium, hard

Return ONLY valid JSON:
{
  "questions": [
    {
      "question": "The question text",
      "category": "financial",
      "difficulty": "medium",
      "suggested_answer": "A strong answer to this question",
      "key_points": ["point1", "point2", "point3"]
    }
  ]
}"""

        slides_context = f"""
Business: {business_idea}
Problem: {', '.join(slides.problem_slide.content)}
Solution: {', '.join(slides.solution_slide.content)}
Market: {', '.join(slides.market_slide.content)}
Team: {', '.join(slides.team_slide.content)}
"""

        user_prompt = f"""Generate {num_questions} investor questions for this pitch:

{slides_context}

Mix of categories (financial, market, team, product, competition).
Mix of difficulties (easy, medium, hard).
Include suggested answers and key points.

Return as JSON."""

        try:
            # Use Llama model for reasoning
            response = await self.groq.generate(
                prompt=user_prompt,
                system_prompt=system_prompt,
                temperature=0.6,
                json_mode=True,
                max_tokens=3000,
            )
            
            data = json.loads(response)
            questions = []
            
            for q_data in data.get("questions", []):
                questions.append(InvestorQuestion(
                    question=q_data.get("question", ""),
                    category=q_data.get("category", "general"),
                    difficulty=q_data.get("difficulty", "medium"),
                    suggested_answer=q_data.get("suggested_answer", ""),
                    key_points=q_data.get("key_points", [])
                ))
            
            return questions
        
        except Exception as e:
            logger.error(f"Error generating investor Q&A: {str(e)}")
            return self._get_default_qa()
    
    # ========================================================================
    # MAIN ORCHESTRATION METHOD
    # ========================================================================
    
    async def create_pitch_deck(
        self,
        business_idea: str,
        business_name: str = "",
        target_market: str = "",
        funding_ask: float = 0,
        brand_tone: str = "professional",
        include_voiceover: bool = True,
        include_demo_script: bool = True,
        include_qa: bool = True
    ) -> PitchDeckResponse:
        """
        Create a complete pitch deck with all features
        """
        
        deck_id = str(uuid.uuid4())
        logger.info(f"Creating pitch deck {deck_id} for: {business_name or business_idea[:50]}")
        
        # Step 1: Generate slides
        logger.info("Generating slides...")
        slides = await self.generate_slides(
            business_idea=business_idea,
            business_name=business_name,
            target_market=target_market,
            funding_ask=funding_ask
        )
        
        # Step 2: Select design theme
        logger.info("Selecting design theme...")
        design_theme = await self.select_design_theme(
            business_idea=business_idea,
            brand_tone=brand_tone
        )
        
        # Step 3: Add charts
        logger.info("Adding charts...")
        slides = await self.add_charts_to_slides(slides, business_idea)
        
        # Step 4: Generate voiceovers (optional)
        voiceovers = []
        if include_voiceover:
            logger.info("Generating voiceovers...")
            voiceovers = await self.generate_voiceovers(slides)
        
        # Step 5: Generate demo script (optional)
        demo_script = None
        if include_demo_script:
            logger.info("Generating demo script...")
            demo_script = await self.generate_demo_script(slides)
        else:
            demo_script = self._get_default_demo_script()
        
        # Step 6: Generate investor Q&A (optional)
        investor_qa = []
        if include_qa:
            logger.info("Generating investor Q&A...")
            investor_qa = await self.generate_investor_qa(business_idea, slides)
        
        # Create response
        response = PitchDeckResponse(
            deck_id=deck_id,
            business_name=business_name or "Your Business",
            tagline=slides.title_slide.content[0] if slides.title_slide.content else "",
            slides=slides,
            design_theme=design_theme,
            voiceovers=voiceovers,
            demo_script=demo_script,
            investor_qa=investor_qa,
            created_at=datetime.now().isoformat()
        )
        
        # Step 7: Export to PPTX (optional)
        if PPTX_AVAILABLE:
            try:
                logger.info("Exporting to PPTX...")
                pptx_bytes = await self.export_to_pptx(response)
                # Removed pptx_bytes from response
                logger.info("PPTX exported successfully")
            except Exception as e:
                logger.warning(f"PPTX export failed: {str(e)}")
        
        logger.info(f"Pitch deck {deck_id} created successfully")
        return response
    
    # ========================================================================
    # HELPER METHODS
    # ========================================================================
    
    def _default_slide(self, slide_number: int, title: str) -> SlideContent:
        """Create a default slide"""
        return SlideContent(
            slide_number=slide_number,
            title=title,
            content=["Content for " + title],
            notes="Notes for " + title
        )
    
    def _get_default_slides(self) -> PitchDeckSlides:
        """Get default slides structure"""
        return PitchDeckSlides(
            title_slide=self._default_slide(1, "Company Name"),
            problem_slide=self._default_slide(2, "The Problem"),
            solution_slide=self._default_slide(3, "Our Solution"),
            market_slide=self._default_slide(4, "Market Opportunity"),
            product_slide=self._default_slide(5, "Product"),
            business_model_slide=self._default_slide(6, "Business Model"),
            traction_slide=self._default_slide(7, "Traction"),
            competition_slide=self._default_slide(8, "Competition"),
            team_slide=self._default_slide(9, "Our Team"),
            financials_slide=self._default_slide(10, "Financials"),
            ask_slide=self._default_slide(11, "The Ask"),
            closing_slide=self._default_slide(12, "Thank You")
        )
    
    def _get_default_theme(self) -> DesignTheme:
        """Get default design theme"""
        return DesignTheme(
            name="Professional",
            primary_color="#2563eb",
            secondary_color="#1e40af",
            accent_color="#f59e0b",
            background_color="#ffffff",
            text_color="#1f2937",
            font_family="Inter, sans-serif",
            style_description="Clean and professional design"
        )
    
    def _get_default_demo_script(self) -> DemoScript:
        """Get default demo script"""
        return DemoScript(
            full_script="Your pitch script will be generated here.",
            slide_scripts=[],
            total_duration_minutes=5.0,
            pacing_cues=["Speak clearly", "Make eye contact"],
            emphasis_points=["Problem", "Solution", "Ask"]
        )
    
    def _get_default_qa(self) -> List[InvestorQuestion]:
        """Get default Q&A questions"""
        return [
            InvestorQuestion(
                question="What is your customer acquisition strategy?",
                category="market",
                difficulty="medium",
                suggested_answer="We plan to acquire customers through...",
                key_points=["Channel strategy", "CAC", "Growth plan"]
            )
        ]
    
    # ========================================================================
    # MODULE 7: SLIDE EXPORT AS PPTX
    # ========================================================================
    
    async def export_to_pptx(
        self,
        deck: PitchDeckResponse,
    ) -> bytes:
        """
        Export pitch deck to PowerPoint format
        """
        
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not installed. Cannot export to PPTX.")
            raise ImportError("python-pptx is required for PPTX export. Install with: pip install python-pptx")
        
        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Parse colors from hex
            primary_color = self._hex_to_rgb(deck.design_theme.primary_color)
            secondary_color = self._hex_to_rgb(deck.design_theme.secondary_color)
            accent_color = self._hex_to_rgb(deck.design_theme.accent_color)
            text_color = self._hex_to_rgb(deck.design_theme.text_color)
            
            # Get all slides
            all_slides = [
                deck.slides.title_slide,
                deck.slides.problem_slide,
                deck.slides.solution_slide,
                deck.slides.market_slide,
                deck.slides.product_slide,
                deck.slides.business_model_slide,
                deck.slides.traction_slide,
                deck.slides.competition_slide,
                deck.slides.team_slide,
                deck.slides.financials_slide,
                deck.slides.ask_slide,
                deck.slides.closing_slide
            ]
            
            # Create each slide
            for slide_content in all_slides:
                if slide_content.slide_number == 1:
                    # Title slide
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    self._create_title_slide(slide, deck, primary_color, text_color)
                else:
                    # Content slide
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    self._create_content_slide(
                        slide, 
                        slide_content, 
                        primary_color, 
                        secondary_color,
                        accent_color,
                        text_color
                    )
                    
                    # Add chart if available
                    if slide_content.chart_data and slide_content.chart_type:
                        await self._add_chart_to_slide(slide, slide_content)
            
            # Save presentation to bytes in memory
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            logger.info("PPTX exported successfully")
            
            return buffer.getvalue()
        
        except Exception as e:
            logger.error(f"Error exporting to PPTX: {str(e)}")
            raise
    
    def _create_title_slide(
        self,
        slide,
        deck: PitchDeckResponse,
        primary_color: RGBColor,
        text_color: RGBColor
    ):
        """Create title slide"""
        
        # Add background rectangle
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = primary_color
        background.line.fill.background()
        
        # Add company name
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = deck.business_name
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add tagline
        tagline_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(8), Inches(0.8)
        )
        tagline_frame = tagline_box.text_frame
        tagline_frame.text = deck.tagline
        tagline_para = tagline_frame.paragraphs[0]
        tagline_para.font.size = Pt(24)
        tagline_para.font.color.rgb = RGBColor(255, 255, 255)
        tagline_para.alignment = PP_ALIGN.CENTER
    
    def _create_content_slide(
        self,
        slide,
        slide_content: SlideContent,
        primary_color: RGBColor,
        secondary_color: RGBColor,
        accent_color: RGBColor,
        text_color: RGBColor
    ):
        """Create content slide"""
        
        # Add header bar
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = primary_color
        header.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(9), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_content.title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add content bullets
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(8.4), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, bullet_text in enumerate(slide_content.content):
            if i > 0:
                content_frame.add_paragraph()
            
            para = content_frame.paragraphs[i]
            para.text = bullet_text
            para.level = 0
            para.font.size = Pt(20)
            para.font.color.rgb = text_color
            para.space_before = Pt(12)
            para.space_after = Pt(12)
        
        # Add slide number
        slide_num_box = slide.shapes.add_textbox(
            Inches(9), Inches(7),
            Inches(0.5), Inches(0.3)
        )
        slide_num_frame = slide_num_box.text_frame
        slide_num_frame.text = str(slide_content.slide_number)
        slide_num_para = slide_num_frame.paragraphs[0]
        slide_num_para.font.size = Pt(14)
        slide_num_para.font.color.rgb = text_color
        slide_num_para.alignment = PP_ALIGN.RIGHT
    
    async def _add_chart_to_slide(self, slide, slide_content: SlideContent):
        """Add chart image to slide"""
        
        try:
            # Generate chart image
            chart_image = await self.quickchart.generate_chart_image(
                chart_type=slide_content.chart_type,
                data=slide_content.chart_data,
                width=600,
                height=400
            )
            
            if chart_image:
                # Save to temporary file
                temp_chart_path = f"temp_chart_{slide_content.slide_number}.png"
                with open(temp_chart_path, 'wb') as f:
                    f.write(chart_image)
                
                # Add to slide
                slide.shapes.add_picture(
                    temp_chart_path,
                    Inches(5), Inches(2),
                    width=Inches(4), height=Inches(3)
                )
                
                # Clean up temp file
                import os
                os.remove(temp_chart_path)
        
        except Exception as e:
            logger.warning(f"Could not add chart to slide {slide_content.slide_number}: {str(e)}")
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        
        hex_color = hex_color.lstrip('#')
        
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        except:
            # Default to blue if parsing fails
            return RGBColor(37, 99, 235)


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def format_currency(amount: float) -> str:
    """Format currency with commas and dollar sign"""
    return f"${amount:,.0f}"


def calculate_slide_duration(word_count: int, words_per_minute: int = 150) -> float:
    """Calculate estimated slide duration in seconds"""
    return (word_count / words_per_minute) * 60


# ============================================================================
# MAIN EXECUTION
# ============================================================================

async def main():
    """Test the Pitch Deck Agent"""
    
    print("=" * 80)
    print("PITCH DECK AGENT - Test")
    print("=" * 80)
    
    try:
        agent = PitchDeckAgent()
        print("✓ Agent initialized successfully\n")
        
        # Test pitch deck creation
        business_idea = "An AI-powered platform that helps startups create investor-ready pitch decks automatically"
        business_name = "DeckAI"
        
        print(f"Creating pitch deck for: {business_name}")
        print(f"Idea: {business_idea}\n")
        
        deck = await agent.create_pitch_deck(
            business_idea=business_idea,
            business_name=business_name,
            target_market="Early-stage startups and entrepreneurs",
            funding_ask=500000,
            brand_tone="modern",
            include_voiceover=False,  # Skip voiceover for quick test
            include_demo_script=True,
            include_qa=True
        )
        
        print("\n" + "=" * 80)
        print("PITCH DECK CREATED SUCCESSFULLY")
        print("=" * 80)
        print(f"\nDeck ID: {deck.deck_id}")
        print(f"Business: {deck.business_name}")
        print(f"Tagline: {deck.tagline}")
        print(f"\nDesign Theme: {deck.design_theme.name}")
        print(f"Colors: {deck.design_theme.primary_color}, {deck.design_theme.secondary_color}")
        
        print(f"\n📊 SLIDES GENERATED:")
        all_slides = [
            deck.slides.title_slide, deck.slides.problem_slide, deck.slides.solution_slide,
            deck.slides.market_slide, deck.slides.product_slide, deck.slides.business_model_slide,
            deck.slides.traction_slide, deck.slides.competition_slide, deck.slides.team_slide,
            deck.slides.financials_slide, deck.slides.ask_slide, deck.slides.closing_slide
        ]
        
        for slide in all_slides:
            print(f"\n  Slide {slide.slide_number}: {slide.title}")
            for item in slide.content[:2]:  # Show first 2 bullet points
                print(f"    • {item}")
            if slide.chart_type:
                print(f"    📈 Chart: {slide.chart_type}")
        
        print(f"\n🎤 DEMO SCRIPT:")
        print(f"  Duration: {deck.demo_script.total_duration_minutes} minutes")
        print(f"  Pacing cues: {len(deck.demo_script.pacing_cues)}")
        
        print(f"\n❓ INVESTOR Q&A:")
        print(f"  Questions generated: {len(deck.investor_qa)}")
        for i, qa in enumerate(deck.investor_qa[:3], 1):
            print(f"\n  Q{i} [{qa.category}] ({qa.difficulty}): {qa.question}")
        
        print("\n" + "=" * 80)
        print("✓ Test completed successfully!")
        print("=" * 80)
        
    except Exception as e:
        print(f"\n✗ Error: {str(e)}")
        logger.exception("Error in main execution")


if __name__ == "__main__":
    asyncio.run(main())
